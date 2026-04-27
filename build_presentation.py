"""
Génère la présentation PowerPoint orientée résultat pour Atlantic Re.
Cible : Direction Générale & Comité de Pilotage.
Palette : Navy / Olive (charte de l'application).
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn
from copy import deepcopy
from lxml import etree


# ── Palette dérivée de tailwind.config.js ────────────────────────────────────
NAVY        = RGBColor(0x2C, 0x47, 0x62)   # primaire
NAVY_DARK   = RGBColor(0x1C, 0x2F, 0x46)
NAVY_900    = RGBColor(0x0F, 0x1C, 0x2D)
NAVY_LIGHT  = RGBColor(0x3D, 0x5C, 0x7C)
OLIVE       = RGBColor(0x73, 0x8C, 0x2C)   # accent
OLIVE_LIGHT = RGBColor(0x92, 0xAF, 0x38)
OLIVE_PALE  = RGBColor(0xE6, 0xEE, 0xCC)
SURFACE     = RGBColor(0xF5, 0xF6, 0xF8)
SURFACE_2   = RGBColor(0xEC, 0xEF, 0xF3)
WHITE       = RGBColor(0xFF, 0xFF, 0xFF)
GRAY_300    = RGBColor(0xC9, 0xCE, 0xD6)
GRAY_500    = RGBColor(0x80, 0x88, 0x96)
GRAY_700    = RGBColor(0x40, 0x4B, 0x5B)
GRAY_900    = RGBColor(0x1F, 0x26, 0x33)
DANGER      = RGBColor(0xDA, 0x3D, 0x44)
WARNING     = RGBColor(0xF6, 0xA8, 0x56)
AMBER       = RGBColor(0xFC, 0xBE, 0x25)


# ── Slide 16:9 ───────────────────────────────────────────────────────────────
prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)

SW = prs.slide_width
SH = prs.slide_height

BLANK = prs.slide_layouts[6]


# ─────────────────────────────────────────────────────────────────────────────
#  Helpers de construction
# ─────────────────────────────────────────────────────────────────────────────
def add_rect(slide, x, y, w, h, fill, line_color=None, line_width=0):
    shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill
    if line_color is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line_color
        shp.line.width = Pt(line_width)
    shp.shadow.inherit = False
    return shp


def add_round_rect(slide, x, y, w, h, fill, line_color=None, line_width=0, corner=0.06):
    shp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    # adjuster (corner radius) — 0..0.5
    shp.adjustments[0] = corner
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill
    if line_color is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line_color
        shp.line.width = Pt(line_width)
    shp.shadow.inherit = False
    return shp


def add_text(slide, x, y, w, h, text, *,
             font="Calibri", size=14, bold=False, italic=False,
             color=GRAY_900, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
             line_spacing=1.15):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.margin_left = tf.margin_right = Emu(0)
    tf.margin_top = tf.margin_bottom = Emu(0)
    tf.word_wrap = True
    tf.vertical_anchor = anchor

    if isinstance(text, str):
        text = [text]

    for i, t in enumerate(text):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.line_spacing = line_spacing
        run = p.add_run()
        run.text = t
        run.font.name = font
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.italic = italic
        run.font.color.rgb = color
    return tb


def add_bullets(slide, x, y, w, h, items, *,
                font="Calibri", size=13, color=GRAY_700, bullet_color=OLIVE,
                line_spacing=1.30, marker="●"):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.margin_left = tf.margin_right = Emu(0)
    tf.margin_top = tf.margin_bottom = Emu(0)
    tf.word_wrap = True
    for i, it in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        p.line_spacing = line_spacing
        p.space_after = Pt(4)
        # marker run
        r1 = p.add_run()
        r1.text = f"{marker}  "
        r1.font.name = font
        r1.font.size = Pt(size)
        r1.font.bold = True
        r1.font.color.rgb = bullet_color
        # text run
        r2 = p.add_run()
        r2.text = it
        r2.font.name = font
        r2.font.size = Pt(size)
        r2.font.color.rgb = color
    return tb


def page_background(slide, fill=SURFACE):
    add_rect(slide, 0, 0, SW, SH, fill)


def page_header(slide, eyebrow, title, *, accent=OLIVE, on_dark=False):
    """Titre de section pour les slides intérieures."""
    bar_color = OLIVE
    add_rect(slide, 0, 0, SW, Inches(0.55), NAVY)
    # accent stripe
    add_rect(slide, 0, Inches(0.55), SW, Inches(0.06), bar_color)
    # eyebrow (small, top-left)
    add_text(slide, Inches(0.55), Inches(0.10), Inches(8), Inches(0.35),
             eyebrow.upper(), size=10, bold=True, color=OLIVE_LIGHT,
             align=PP_ALIGN.LEFT)
    # right-side branding mark
    add_text(slide, Inches(11.0), Inches(0.10), Inches(2.0), Inches(0.35),
             "ATLANTIC RE", size=10, bold=True, color=WHITE,
             align=PP_ALIGN.RIGHT)
    # main title
    add_text(slide, Inches(0.55), Inches(0.78), Inches(12.2), Inches(0.7),
             title, size=26, bold=True, color=NAVY_DARK, align=PP_ALIGN.LEFT)
    # subtle separator under title
    add_rect(slide, Inches(0.55), Inches(1.50), Inches(0.7), Inches(0.04), OLIVE)


def page_footer(slide, page_num, total, section_label=""):
    add_rect(slide, 0, Inches(7.18), SW, Inches(0.32), NAVY_DARK)
    add_text(slide, Inches(0.55), Inches(7.20), Inches(8), Inches(0.28),
             section_label, size=9, bold=True, color=OLIVE_LIGHT,
             align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.MIDDLE)
    add_text(slide, Inches(11.0), Inches(7.20), Inches(2.0), Inches(0.28),
             f"{page_num:02d} / {total:02d}", size=9, color=WHITE,
             align=PP_ALIGN.RIGHT, anchor=MSO_ANCHOR.MIDDLE)


def add_screenshot_zone(slide, x, y, w, h, label="Capture d'écran"):
    """Rectangle dashé clair – à remplacer par une capture après export."""
    shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    shp.fill.solid()
    shp.fill.fore_color.rgb = WHITE
    shp.line.color.rgb = OLIVE
    shp.line.width = Pt(1.5)
    # dash style (solid line by default — set DASH via XML)
    ln = shp.line._get_or_add_ln()
    prstDash = etree.SubElement(ln, qn('a:prstDash'))
    prstDash.set('val', 'dash')
    shp.shadow.inherit = False

    # corner marks (decorative)
    cm = Inches(0.18)
    for cx, cy in [(x, y), (x + w - cm, y), (x, y + h - cm), (x + w - cm, y + h - cm)]:
        c = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, cx, cy, cm, cm)
        c.fill.solid(); c.fill.fore_color.rgb = OLIVE
        c.line.fill.background(); c.shadow.inherit = False

    # central pictogram (camera-like simple square)
    icon_w = Inches(0.55)
    icon_h = Inches(0.40)
    icon_x = x + w/2 - icon_w/2
    icon_y = y + h/2 - icon_h - Inches(0.30)
    ic = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, icon_x, icon_y, icon_w, icon_h)
    ic.adjustments[0] = 0.20
    ic.fill.solid(); ic.fill.fore_color.rgb = NAVY
    ic.line.fill.background(); ic.shadow.inherit = False
    lens = slide.shapes.add_shape(MSO_SHAPE.OVAL,
                                  icon_x + icon_w/2 - Inches(0.13),
                                  icon_y + icon_h/2 - Inches(0.13),
                                  Inches(0.26), Inches(0.26))
    lens.fill.solid(); lens.fill.fore_color.rgb = OLIVE_LIGHT
    lens.line.color.rgb = WHITE; lens.line.width = Pt(1.2)
    lens.shadow.inherit = False

    # main label
    add_text(slide, x, y + h/2, w, Inches(0.35),
             label, size=12, bold=True, color=NAVY,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    # hint
    add_text(slide, x, y + h/2 + Inches(0.30), w, Inches(0.30),
             "Zone réservée — insérer la capture après export",
             size=9, italic=True, color=GRAY_500,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)


def kpi_card(slide, x, y, w, h, value, label, *, accent=OLIVE, sub=None):
    add_round_rect(slide, x, y, w, h, WHITE, line_color=GRAY_300, line_width=0.75, corner=0.10)
    # left accent bar
    add_rect(slide, x, y + Inches(0.25), Inches(0.10), h - Inches(0.5), accent)
    add_text(slide, x + Inches(0.30), y + Inches(0.20), w - Inches(0.40), Inches(0.65),
             value, size=26, bold=True, color=NAVY_DARK, align=PP_ALIGN.LEFT)
    add_text(slide, x + Inches(0.30), y + Inches(0.85), w - Inches(0.40), Inches(0.40),
             label, size=11, bold=True, color=GRAY_700, align=PP_ALIGN.LEFT)
    if sub:
        add_text(slide, x + Inches(0.30), y + Inches(1.18), w - Inches(0.40), Inches(0.30),
                 sub, size=9, italic=True, color=GRAY_500, align=PP_ALIGN.LEFT)


def section_divider(eyebrow, title, subtitle, page_num, total):
    """Slide intercalaire pleine page, fond sombre."""
    s = prs.slides.add_slide(BLANK)
    add_rect(s, 0, 0, SW, SH, NAVY_DARK)
    # decorative diagonal block
    block = s.shapes.add_shape(MSO_SHAPE.PARALLELOGRAM,
                               Inches(8.5), Inches(-0.8), Inches(7), Inches(9.3))
    block.adjustments[0] = 0.45
    block.fill.solid(); block.fill.fore_color.rgb = NAVY
    block.line.fill.background(); block.shadow.inherit = False
    # olive bar
    add_rect(s, 0, Inches(2.4), Inches(1.0), Inches(0.10), OLIVE)
    add_text(s, Inches(0.55), Inches(2.5), Inches(8), Inches(0.5),
             eyebrow.upper(), size=14, bold=True, color=OLIVE_LIGHT)
    add_text(s, Inches(0.55), Inches(3.0), Inches(11), Inches(1.6),
             title, size=44, bold=True, color=WHITE, line_spacing=1.05)
    add_text(s, Inches(0.55), Inches(4.7), Inches(11), Inches(1.0),
             subtitle, size=16, italic=True, color=GRAY_300, line_spacing=1.30)
    # bottom right page
    add_text(s, Inches(11.0), Inches(7.0), Inches(2), Inches(0.4),
             f"{page_num:02d} / {total:02d}", size=10, color=OLIVE_LIGHT,
             align=PP_ALIGN.RIGHT)
    return s


def pill(slide, x, y, w, h, text, fill, text_color=WHITE, *, font_size=10, bold=True):
    p = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    p.adjustments[0] = 0.50
    p.fill.solid(); p.fill.fore_color.rgb = fill
    p.line.fill.background(); p.shadow.inherit = False
    tf = p.text_frame
    tf.margin_left = Emu(0); tf.margin_right = Emu(0)
    tf.margin_top = Emu(0); tf.margin_bottom = Emu(0)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    pp = tf.paragraphs[0]
    pp.alignment = PP_ALIGN.CENTER
    r = pp.add_run()
    r.text = text
    r.font.name = "Calibri"; r.font.size = Pt(font_size)
    r.font.bold = bold; r.font.color.rgb = text_color
    return p


# ─────────────────────────────────────────────────────────────────────────────
#  Plan : on construit d'abord la liste pour calculer le total avant rendu
# ─────────────────────────────────────────────────────────────────────────────
TOTAL = 36


# =============================================================================
#  SLIDE 1 — COUVERTURE
# =============================================================================
def slide_cover():
    s = prs.slides.add_slide(BLANK)
    add_rect(s, 0, 0, SW, SH, NAVY_DARK)
    # large diagonal block right
    block = s.shapes.add_shape(MSO_SHAPE.PARALLELOGRAM,
                               Inches(7.5), Inches(-1), Inches(8), Inches(10))
    block.adjustments[0] = 0.42
    block.fill.solid(); block.fill.fore_color.rgb = NAVY
    block.line.fill.background(); block.shadow.inherit = False

    # olive accent
    add_rect(s, 0, Inches(2.3), Inches(0.85), Inches(0.10), OLIVE)

    # eyebrow
    add_text(s, Inches(0.55), Inches(2.45), Inches(8), Inches(0.5),
             "PROJET DE FIN D'ÉTUDES — ATLANTIC RE", size=13,
             bold=True, color=OLIVE_LIGHT)
    # main title
    add_text(s, Inches(0.55), Inches(2.95), Inches(11), Inches(1.7),
             "Plate-forme intégrée de\ngouvernance et d'aide à la décision",
             size=40, bold=True, color=WHITE, line_spacing=1.05)
    # subtitle
    add_text(s, Inches(0.55), Inches(4.85), Inches(11), Inches(0.9),
             "Pipeline ETL  ·  Plate-forme métier  ·  Scoring multicritère pour Reach2030",
             size=16, italic=True, color=GRAY_300)

    # bottom band
    add_rect(s, 0, Inches(6.85), SW, Inches(0.65), NAVY_900)
    add_text(s, Inches(0.55), Inches(6.92), Inches(8), Inches(0.5),
             "Présentation Direction Générale  ·  Exercice 2026",
             size=12, bold=True, color=WHITE, anchor=MSO_ANCHOR.MIDDLE)
    add_text(s, Inches(8), Inches(6.92), Inches(5), Inches(0.5),
             "Pôle Business Développement",
             size=11, color=OLIVE_LIGHT, align=PP_ALIGN.RIGHT,
             anchor=MSO_ANCHOR.MIDDLE)

    # logo placeholder square (top-left)
    add_round_rect(s, Inches(0.55), Inches(0.55), Inches(2.2), Inches(0.65),
                   OLIVE, corner=0.30)
    add_text(s, Inches(0.55), Inches(0.55), Inches(2.2), Inches(0.65),
             "ATLANTIC RE", size=14, bold=True, color=WHITE,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)


# =============================================================================
#  SLIDE 2 — AGENDA
# =============================================================================
def slide_agenda():
    s = prs.slides.add_slide(BLANK)
    page_background(s)
    page_header(s, "Sommaire exécutif", "Agenda de la présentation")

    items = [
        ("01", "Contexte stratégique",        "Reach2030 et ambition continentale"),
        ("02", "Benchmarking sectoriel",      "Marché africain et paysage concurrentiel"),
        ("03", "Analyse des besoins",         "Démarche itérative et constat terrain"),
        ("04", "La solution Atlantic Re",     "Architecture et trajectoire de livraison"),
        ("05", "Axe 1 — Pipeline ETL",        "Extraction, transformation, chargement"),
        ("06", "Axe 2 — Plate-forme métier",  "Restitution, analyses, gouvernance"),
        ("07", "Axe 3 — Machine Learning",    "Cartographies, scoring SCAR, recommandations"),
        ("08", "Convergence & bilan",         "Synergie, indicateurs et perspectives"),
    ]

    col_w = Inches(6.0)
    row_h = Inches(0.62)
    gap   = Inches(0.10)
    start_y = Inches(1.95)

    for i, (num, title, desc) in enumerate(items):
        col = i // 4
        row = i % 4
        x = Inches(0.55) + col * (col_w + Inches(0.30))
        y = start_y + row * (row_h + gap)

        add_round_rect(s, x, y, col_w, row_h, WHITE,
                       line_color=GRAY_300, line_width=0.5, corner=0.18)
        add_round_rect(s, x, y, Inches(0.62), row_h, NAVY, corner=0.18)
        add_text(s, x, y, Inches(0.62), row_h, num, size=18, bold=True,
                 color=OLIVE_LIGHT, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        add_text(s, x + Inches(0.78), y + Inches(0.05), col_w - Inches(0.85), Inches(0.30),
                 title, size=13, bold=True, color=NAVY_DARK)
        add_text(s, x + Inches(0.78), y + Inches(0.32), col_w - Inches(0.85), Inches(0.28),
                 desc, size=10, italic=True, color=GRAY_500)

    page_footer(s, 2, TOTAL, "Sommaire exécutif")


# =============================================================================
#  SECTION 1 — CONTEXTE STRATÉGIQUE
# =============================================================================
def slide_section1_divider():
    section_divider("Section 01",
                    "Contexte stratégique et\nambition continentale",
                    "Le plan Reach2030 et le rôle du Pôle Business Développement.",
                    3, TOTAL)


def slide_reach2030():
    s = prs.slides.add_slide(BLANK)
    page_background(s)
    page_header(s, "01 · Contexte stratégique",
                "Reach2030 — l'Afrique comme axe prioritaire de croissance")

    # Left column — texte
    add_text(s, Inches(0.55), Inches(1.85), Inches(7.5), Inches(0.4),
             "L'opportunité africaine en chiffres", size=14, bold=True, color=NAVY_DARK)
    add_bullets(s, Inches(0.55), Inches(2.30), Inches(7.4), Inches(3.5),
        [
            "Pénétration de l'assurance < 3 % du PIB en Afrique subsaharienne (vs >7 % au niveau mondial).",
            "Croissance démographique soutenue, urbanisation accélérée, projets d'infrastructure massifs.",
            "Exposition climatique croissante — demande structurelle de couverture en hausse.",
            "Capacité locale de souscription insuffisante : appel structurel à la réassurance internationale.",
            "Atlantic Re bénéficie d'un ancrage marocain neutre et du soutien institutionnel du Groupe CDG.",
        ],
        size=12)

    # Right column — KPIs Atlantic Re 2025
    add_text(s, Inches(8.30), Inches(1.85), Inches(4.5), Inches(0.4),
             "Atlantic Re — Exercice 2025", size=14, bold=True, color=NAVY_DARK)

    kpi_card(s, Inches(8.30), Inches(2.30), Inches(2.20), Inches(1.40),
             "4 039\nMDH", "Chiffre d'affaires", accent=OLIVE)
    kpi_card(s, Inches(10.60), Inches(2.30), Inches(2.20), Inches(1.40),
             "419\nMDH", "Résultat net", accent=OLIVE)
    kpi_card(s, Inches(8.30), Inches(3.85), Inches(2.20), Inches(1.40),
             "92,3 %", "Ratio combiné", accent=NAVY_LIGHT)
    kpi_card(s, Inches(10.60), Inches(3.85), Inches(2.20), Inches(1.40),
             "215 %", "Ratio S2", accent=NAVY_LIGHT, sub="+9 pts vs N-1")
    kpi_card(s, Inches(8.30), Inches(5.40), Inches(2.20), Inches(1.40),
             "+ 60", "Pays couverts", accent=OLIVE_LIGHT)
    kpi_card(s, Inches(10.60), Inches(5.40), Inches(2.20), Inches(1.40),
             "+ 400", "Clients actifs", accent=OLIVE_LIGHT)

    # Quote band en bas gauche
    add_round_rect(s, Inches(0.55), Inches(6.05), Inches(7.4), Inches(1.0),
                   NAVY, corner=0.10)
    add_text(s, Inches(0.85), Inches(6.15), Inches(7.0), Inches(0.30),
             "ENJEU DIRECTION GÉNÉRALE", size=9, bold=True, color=OLIVE_LIGHT)
    add_text(s, Inches(0.85), Inches(6.42), Inches(7.0), Inches(0.55),
             "Identifier — parmi des marchés africains très hétérogènes — ceux qui offrent\nla meilleure combinaison potentiel × rentabilité × stabilité × dynamisme.",
             size=11, color=WHITE, italic=True, line_spacing=1.20)

    page_footer(s, 4, TOTAL, "01 · Contexte stratégique")


def slide_pole_bd():
    s = prs.slides.add_slide(BLANK)
    page_background(s)
    page_header(s, "01 · Contexte stratégique",
                "Le Pôle Business Développement — cadre et défis")

    # left : mission
    add_round_rect(s, Inches(0.55), Inches(1.85), Inches(6.0), Inches(2.55),
                   WHITE, line_color=GRAY_300, line_width=0.5, corner=0.06)
    add_rect(s, Inches(0.55), Inches(1.85), Inches(6.0), Inches(0.45), NAVY)
    add_text(s, Inches(0.75), Inches(1.85), Inches(5.7), Inches(0.45),
             "MISSION DU PÔLE", size=11, bold=True, color=WHITE,
             anchor=MSO_ANCHOR.MIDDLE)
    add_text(s, Inches(0.75), Inches(2.45), Inches(5.6), Inches(0.5),
             "Identifier et analyser les opportunités de développement",
             size=14, bold=True, color=NAVY_DARK)
    add_bullets(s, Inches(0.75), Inches(3.05), Inches(5.5), Inches(1.45),
        [
            "Alimentation des décisions de souscription Direction Générale.",
            "Études de marché et veille concurrentielle continentale.",
            "Mesure de l'attractivité des marchés cibles (potentiel, risque).",
        ], size=11)

    # right : défis
    add_round_rect(s, Inches(6.85), Inches(1.85), Inches(6.0), Inches(2.55),
                   WHITE, line_color=GRAY_300, line_width=0.5, corner=0.06)
    add_rect(s, Inches(6.85), Inches(1.85), Inches(6.0), Inches(0.45), DANGER)
    add_text(s, Inches(7.05), Inches(1.85), Inches(5.7), Inches(0.45),
             "DÉFIS RENCONTRÉS", size=11, bold=True, color=WHITE,
             anchor=MSO_ANCHOR.MIDDLE)
    add_text(s, Inches(7.05), Inches(2.45), Inches(5.6), Inches(0.5),
             "Outils manquants pour exploiter la donnée",
             size=14, bold=True, color=NAVY_DARK)
    add_bullets(s, Inches(7.05), Inches(3.05), Inches(5.5), Inches(1.45),
        [
            "Aucun portail de consultation des KPIs internes.",
            "Manipulations Excel manuelles, versions divergentes.",
            "Aucun référentiel externe consolidé pour benchmarker.",
        ], size=11, bullet_color=DANGER)

    # bottom timeline
    add_text(s, Inches(0.55), Inches(4.70), Inches(12.2), Inches(0.4),
             "Cadre d'intervention — 4 mois en immersion (Février → Juin 2026)",
             size=14, bold=True, color=NAVY_DARK)

    steps = [
        ("S0", "Immersion\n& cadrage"),
        ("S1", "Analyse\nbesoins"),
        ("S2", "Architecture"),
        ("S3-4", "Portail\nAxe 1"),
        ("S5", "Données ext.\nAxe 2"),
        ("S6", "Modèle\nSCAR"),
        ("S7", "Clôture"),
    ]
    n = len(steps)
    avail = Inches(12.2)
    step_w = avail / n
    yT = Inches(5.40)
    for i, (lbl, txt) in enumerate(steps):
        x = Inches(0.55) + step_w * i
        # circle
        cx = x + step_w/2 - Inches(0.30)
        circle = s.shapes.add_shape(MSO_SHAPE.OVAL, cx, yT, Inches(0.60), Inches(0.60))
        circle.fill.solid()
        circle.fill.fore_color.rgb = OLIVE if i < n-1 else NAVY
        circle.line.fill.background(); circle.shadow.inherit = False
        add_text(s, cx, yT, Inches(0.60), Inches(0.60), lbl,
                 size=11, bold=True, color=WHITE,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        # connector
        if i < n-1:
            add_rect(s, cx + Inches(0.60), yT + Inches(0.27),
                     step_w - Inches(0.60), Inches(0.06), GRAY_300)
        # label
        add_text(s, x, yT + Inches(0.70), step_w, Inches(0.60), txt,
                 size=10, color=GRAY_700, align=PP_ALIGN.CENTER)

    page_footer(s, 5, TOTAL, "01 · Contexte stratégique")


# =============================================================================
#  SECTION 2 — BENCHMARKING SECTORIEL
# =============================================================================
def slide_section2_divider():
    section_divider("Section 02",
                    "Benchmarking sectoriel",
                    "Lecture du marché africain et cartographie des acteurs en présence.",
                    6, TOTAL)


def slide_marche_africain():
    s = prs.slides.add_slide(BLANK)
    page_background(s)
    page_header(s, "02 · Benchmarking sectoriel",
                "Le marché africain — une opportunité structurelle")

    # 3 colonnes : moteurs / freins / implications
    col_y = Inches(1.95)
    col_h = Inches(4.0)
    col_w = Inches(4.05)

    cols = [
        ("MOTEURS", OLIVE,
         "Une demande qui s'amplifie",
         [
             "Croissance démographique +2,5 % / an",
             "Urbanisation et classe moyenne en hausse",
             "Programmes d'infrastructure publics",
             "Multiplication des risques climatiques",
             "ZLECAf — intégration commerciale",
         ]),
        ("FREINS", DANGER,
         "Asymétrie capacité / besoins",
         [
             "Pénétration < 3 % PIB (vs >7 % monde)",
             "Capacités locales souvent fragmentées",
             "Hétérogénéité réglementaire (CIMA, etc.)",
             "Volatilité macroéconomique localisée",
             "Provisionnement parfois peu fiable",
         ]),
        ("IMPLICATIONS", NAVY_LIGHT,
         "Pour un réassureur international",
         [
             "Gisement structurel de primes à capter",
             "Nécessité d'une lecture marché par marché",
             "Discipline technique exigée (ULR, S/P)",
             "Positionnement neutre = avantage CDG",
             "Décision stratégique = besoin d'un outil",
         ]),
    ]

    for i, (lbl, color, head, items) in enumerate(cols):
        x = Inches(0.55) + i * (col_w + Inches(0.20))
        add_round_rect(s, x, col_y, col_w, col_h,
                       WHITE, line_color=GRAY_300, line_width=0.5, corner=0.06)
        add_rect(s, x, col_y, col_w, Inches(0.40), color)
        add_text(s, x + Inches(0.20), col_y, col_w - Inches(0.20), Inches(0.40),
                 lbl, size=10, bold=True, color=WHITE, anchor=MSO_ANCHOR.MIDDLE)
        add_text(s, x + Inches(0.20), col_y + Inches(0.55),
                 col_w - Inches(0.40), Inches(0.45),
                 head, size=14, bold=True, color=NAVY_DARK)
        add_bullets(s, x + Inches(0.20), col_y + Inches(1.05),
                    col_w - Inches(0.40), col_h - Inches(1.15),
                    items, size=11, bullet_color=color, line_spacing=1.30)

    # bottom band — synthèse
    add_round_rect(s, Inches(0.55), Inches(6.15), Inches(12.20), Inches(0.85),
                   NAVY, corner=0.10)
    add_text(s, Inches(0.85), Inches(6.20), Inches(11.7), Inches(0.30),
             "INSIGHT CLÉ", size=10, bold=True, color=OLIVE_LIGHT)
    add_text(s, Inches(0.85), Inches(6.45), Inches(11.7), Inches(0.50),
             "L'Afrique est la dernière grande frontière de croissance — la décision de souscription ne peut\nplus reposer sur l'intuition. Elle exige un cadre quantitatif rigoureux et auditable.",
             size=11.5, color=WHITE, italic=True, line_spacing=1.20)

    page_footer(s, 7, TOTAL, "02 · Benchmarking sectoriel")


def slide_paysage_concurrence():
    s = prs.slides.add_slide(BLANK)
    page_background(s)
    page_header(s, "02 · Benchmarking sectoriel",
                "Le paysage concurrentiel africain")

    # 4 catégories
    cat = [
        ("Réassureurs panafricains",      OLIVE,
         ["Africa Re (premier panafricain — Atlantic Re actionnaire)",
          "ZEP-Re (Afrique orientale et australe)",
          "CICA-Re (réassureur institutionnel zone CIMA)"]),
        ("Réassureurs internationaux",    NAVY,
         ["Munich Re, Swiss Re, Hannover Re, Scor",
          "Souscription depuis sièges européens ou bureaux régionaux",
          "Forte capacité de placement, exigence technique"]),
        ("Réassureurs régionaux & locaux", NAVY_LIGHT,
         ["Implantation anglophone, lusophone, CIMA",
          "Proximité réglementaire, relations historiques",
          "Souvent moins capitalisés, périmètres limités"]),
        ("Courtiers globaux",             OLIVE_LIGHT,
         ["Aon, Guy Carpenter, Willis Re",
          "Influence croissante sur la structuration",
          "Accès à la capacité internationale"]),
    ]

    grid_x0 = Inches(0.55)
    grid_y0 = Inches(1.95)
    cw = Inches(6.05)
    ch = Inches(2.20)
    gx = Inches(0.20)
    gy = Inches(0.20)

    for i, (title, color, items) in enumerate(cat):
        col = i % 2
        row = i // 2
        x = grid_x0 + col * (cw + gx)
        y = grid_y0 + row * (ch + gy)
        add_round_rect(s, x, y, cw, ch, WHITE,
                       line_color=GRAY_300, line_width=0.5, corner=0.06)
        # left accent
        add_rect(s, x, y, Inches(0.10), ch, color)
        add_text(s, x + Inches(0.30), y + Inches(0.20), cw - Inches(0.4), Inches(0.40),
                 title, size=14, bold=True, color=NAVY_DARK)
        add_bullets(s, x + Inches(0.30), y + Inches(0.65),
                    cw - Inches(0.5), ch - Inches(0.7),
                    items, size=11, bullet_color=color, line_spacing=1.30)

    # bottom : positionnement Atlantic Re
    add_round_rect(s, Inches(0.55), Inches(6.50), Inches(12.20), Inches(0.55),
                   NAVY_DARK, corner=0.20)
    add_text(s, Inches(0.85), Inches(6.50), Inches(3.8), Inches(0.55),
             "POSITIONNEMENT ATLANTIC RE", size=10, bold=True, color=OLIVE_LIGHT,
             anchor=MSO_ANCHOR.MIDDLE)
    add_text(s, Inches(4.65), Inches(6.50), Inches(8.0), Inches(0.55),
             "Ancrage marocain neutre · 6 décennies d'expertise francophone · Soutien institutionnel CDG",
             size=11, color=WHITE, anchor=MSO_ANCHOR.MIDDLE)

    page_footer(s, 8, TOTAL, "02 · Benchmarking sectoriel")


def slide_benchmark_outils():
    s = prs.slides.add_slide(BLANK)
    page_background(s)
    page_header(s, "02 · Benchmarking sectoriel",
                "Pratiques digitales du marché — où se situer ?")

    # Tableau de benchmark
    rows = [
        ("Pratique",                       "Industrie",                    "Atlantic Re — avant",   "Atlantic Re — visé"),
        ("Plate-forme analytique unifiée", "Standard Tier-1",              "Aucune",                "Plate-forme Atlantic Re"),
        ("Données externes consolidées",   "Achat outils tiers (Axco)",    "Sources éclatées",      "Panel 34 pays × 10 ans"),
        ("Scoring multicritère marché",    "Modèles internes propres",     "Inexistant",            "Modèle SCAR calibré"),
        ("Réconciliation données cédantes","Workflow CRM dédié",           "Manuel, sous Excel",    "Automatisée"),
        ("Sécurité / RBAC / audit",        "Norme ISO 27001",              "Partielle",             "RBAC + journal d'activité"),
        ("Souveraineté de la donnée",      "Hébergement local stratégique","N/A",                   "100 % on-premise"),
    ]

    nx = Inches(0.55)
    ny = Inches(1.95)
    col_widths = [Inches(3.2), Inches(2.9), Inches(2.9), Inches(3.20)]
    rh = Inches(0.55)

    # header
    cx = nx
    for i, w in enumerate(col_widths):
        add_rect(s, cx, ny, w, rh, NAVY)
        add_text(s, cx + Inches(0.15), ny, w - Inches(0.30), rh,
                 rows[0][i], size=11, bold=True, color=WHITE,
                 anchor=MSO_ANCHOR.MIDDLE)
        cx += w
    # body
    for r, row in enumerate(rows[1:], start=1):
        cx = nx
        bg = WHITE if r % 2 == 1 else SURFACE_2
        for i, w in enumerate(col_widths):
            add_rect(s, cx, ny + rh * r, w, rh, bg)
            color = NAVY_DARK
            bold = False
            if i == 0:
                bold = True
            elif i == 2:
                color = DANGER
            elif i == 3:
                color = OLIVE
                bold = True
            add_text(s, cx + Inches(0.15), ny + rh * r, w - Inches(0.30), rh,
                     row[i], size=11, bold=bold, color=color,
                     anchor=MSO_ANCHOR.MIDDLE)
            cx += w

    # bordure inférieure
    add_rect(s, nx, ny + rh * len(rows), sum(col_widths, Emu(0)), Inches(0.04), NAVY)

    page_footer(s, 9, TOTAL, "02 · Benchmarking sectoriel")


# =============================================================================
#  SECTION 3 — ANALYSE DES BESOINS
# =============================================================================
def slide_section3_divider():
    section_divider("Section 03",
                    "Analyse des besoins",
                    "Démarche itérative en trois cycles : du cadrage à la validation MoSCoW.",
                    10, TOTAL)


def slide_demarche_3_cycles():
    s = prs.slides.add_slide(BLANK)
    page_background(s)
    page_header(s, "03 · Analyse des besoins",
                "Une démarche itérative en trois cycles")

    cycles = [
        ("CYCLE 1", "Cadrage initial",
         "Sem. 1-2",
         ["Entretiens semi-directifs encadrant",
          "Étude Reach2030 / Sigma / FANAF",
          "Cartographie des parties prenantes"],
         "Cahier des charges initial",
         NAVY_LIGHT),
        ("CYCLE 2", "Immersion terrain",
         "Sem. 2-3",
         ["Observation participante du Pôle BD",
          "Analyse documents internes",
          "Détection des besoins latents"],
         "Élargissement du périmètre",
         WARNING),
        ("CYCLE 3", "Validation MoSCoW",
         "Sem. 3",
         ["Consolidation BE + BL",
          "Priorisation Must / Should / Could",
          "Validation formelle encadrant"],
         "16 BF + 8 BNF validés",
         OLIVE),
    ]

    cw = Inches(4.05)
    ch = Inches(3.85)
    gap = Inches(0.10)
    y0 = Inches(2.0)

    for i, (head, name, period, items, output, color) in enumerate(cycles):
        x = Inches(0.55) + i * (cw + gap)
        add_round_rect(s, x, y0, cw, ch, WHITE,
                       line_color=GRAY_300, line_width=0.5, corner=0.06)
        # numéro grand
        add_rect(s, x, y0, cw, Inches(0.50), color)
        add_text(s, x + Inches(0.20), y0, cw - Inches(0.40), Inches(0.50),
                 head, size=12, bold=True, color=WHITE, anchor=MSO_ANCHOR.MIDDLE)
        pill(s, x + cw - Inches(1.10), y0 + Inches(0.10),
             Inches(0.95), Inches(0.30), period, NAVY_DARK, font_size=9)
        add_text(s, x + Inches(0.20), y0 + Inches(0.60), cw - Inches(0.40), Inches(0.50),
                 name, size=15, bold=True, color=NAVY_DARK)
        add_bullets(s, x + Inches(0.20), y0 + Inches(1.10),
                    cw - Inches(0.40), Inches(2.0),
                    items, size=11, bullet_color=color, line_spacing=1.35)
        # output band
        add_round_rect(s, x + Inches(0.20), y0 + ch - Inches(0.65),
                       cw - Inches(0.40), Inches(0.50),
                       SURFACE_2, corner=0.30)
        add_text(s, x + Inches(0.20), y0 + ch - Inches(0.65),
                 cw - Inches(0.40), Inches(0.50),
                 f"→  {output}", size=11, bold=True, color=NAVY_DARK,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

        # arrow
        if i < len(cycles) - 1:
            ax = x + cw + Inches(0.005)
            arr = s.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW,
                                     ax - Inches(0.05), y0 + ch/2 - Inches(0.10),
                                     Inches(0.20), Inches(0.20))
            arr.fill.solid(); arr.fill.fore_color.rgb = NAVY
            arr.line.fill.background(); arr.shadow.inherit = False

    # bottom
    add_round_rect(s, Inches(0.55), Inches(6.20), Inches(12.20), Inches(0.80),
                   NAVY, corner=0.10)
    add_text(s, Inches(0.85), Inches(6.25), Inches(11.7), Inches(0.30),
             "PIVOT MÉTHODOLOGIQUE", size=10, bold=True, color=OLIVE_LIGHT)
    add_text(s, Inches(0.85), Inches(6.50), Inches(11.7), Inches(0.45),
             "Le Cycle 2 est décisif : il a transformé un PFE de modélisation pure en un projet de modélisation + digitalisation.",
             size=11.5, color=WHITE, italic=True)

    page_footer(s, 11, TOTAL, "03 · Analyse des besoins")


def slide_constats_terrain():
    s = prs.slides.add_slide(BLANK)
    page_background(s)
    page_header(s, "03 · Analyse des besoins",
                "Le constat terrain — 7 besoins latents identifiés")

    rows = [
        ("BL.1", "Absence de consultation structurée", "Versions divergentes des KPI",        "Critique"),
        ("BL.2", "Absence de filtrage multicritère",  "Tableaux croisés manuels",            "Critique"),
        ("BL.3", "Indicateurs composites non calculés","Recalculs manuels, erreurs",          "Élevée"),
        ("BL.4", "Incohérences de nommage cédantes",  "Agrégations approximatives",          "Élevée"),
        ("BL.5", "Absence de comparaison directe",    "Benchmarking manuel entre marchés",   "Modérée"),
        ("BL.6", "Absence de traçabilité",            "Audit d'usage impossible",            "Modérée"),
        ("BL.7", "Modèle externe non fiabilisable",   "Données internes non assainies",       "Critique"),
    ]
    color_map = {"Critique": DANGER, "Élevée": WARNING, "Modérée": AMBER}

    nx = Inches(0.55)
    ny = Inches(1.95)
    cwd = [Inches(0.85), Inches(4.4), Inches(4.6), Inches(2.35)]
    rh = Inches(0.50)

    headers = ["Réf.", "Besoin identifié", "Impact opérationnel", "Criticité"]
    cx = nx
    for i, w in enumerate(cwd):
        add_rect(s, cx, ny, w, rh, NAVY)
        add_text(s, cx + Inches(0.15), ny, w - Inches(0.30), rh,
                 headers[i], size=11, bold=True, color=WHITE,
                 anchor=MSO_ANCHOR.MIDDLE)
        cx += w

    for r, row in enumerate(rows, start=1):
        cx = nx
        bg = WHITE if r % 2 == 1 else SURFACE_2
        for i, w in enumerate(cwd):
            add_rect(s, cx, ny + rh * r, w, rh, bg)
            if i < 3:
                add_text(s, cx + Inches(0.15), ny + rh * r, w - Inches(0.30), rh,
                         row[i], size=11,
                         bold=(i == 0 or i == 1),
                         color=NAVY_DARK if i == 1 else GRAY_700,
                         anchor=MSO_ANCHOR.MIDDLE)
            else:
                # pill criticité
                pill_w = Inches(1.30); pill_h = Inches(0.32)
                px = cx + (w - pill_w)/2
                py = ny + rh * r + (rh - pill_h)/2
                pill(s, px, py, pill_w, pill_h, row[i], color_map[row[i]],
                     font_size=10)
            cx += w

    # bottom takeaway
    add_round_rect(s, Inches(0.55), Inches(6.10), Inches(12.20), Inches(0.95),
                   NAVY_DARK, corner=0.10)
    add_text(s, Inches(0.85), Inches(6.20), Inches(11.7), Inches(0.30),
             "CONSÉQUENCE STRATÉGIQUE", size=10, bold=True, color=OLIVE_LIGHT)
    add_text(s, Inches(0.85), Inches(6.50), Inches(11.7), Inches(0.50),
             "Sans assainissement préalable de la donnée interne, tout modèle externe d'attractivité de marché\nserait construit sur un socle non fiable. La digitalisation devient un prérequis, pas un complément.",
             size=11.5, color=WHITE, italic=True, line_spacing=1.20)

    page_footer(s, 12, TOTAL, "03 · Analyse des besoins")


def slide_elargissement_perimetre():
    s = prs.slides.add_slide(BLANK)
    page_background(s)
    page_header(s, "03 · Analyse des besoins",
                "Élargissement du périmètre — trois axes complémentaires")

    # AVANT
    add_round_rect(s, Inches(0.55), Inches(1.85), Inches(5.5), Inches(2.05),
                   WHITE, line_color=GRAY_300, line_width=0.5, corner=0.06)
    add_rect(s, Inches(0.55), Inches(1.85), Inches(5.5), Inches(0.40), GRAY_500)
    add_text(s, Inches(0.75), Inches(1.85), Inches(5.0), Inches(0.40),
             "PÉRIMÈTRE INITIAL — COMMANDE", size=10, bold=True, color=WHITE,
             anchor=MSO_ANCHOR.MIDDLE)
    add_text(s, Inches(0.75), Inches(2.40), Inches(5.0), Inches(0.50),
             "Modélisation pure",
             size=17, bold=True, color=NAVY_DARK)
    add_bullets(s, Inches(0.75), Inches(2.95), Inches(5.0), Inches(0.95),
        [
            "Modèle de scoring multicritère.",
            "Données externes longitudinales 10 ans.",
            "Restitution cartographique.",
        ], size=11)

    # flèche
    arr = s.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW,
                             Inches(6.40), Inches(2.65),
                             Inches(0.45), Inches(0.55))
    arr.fill.solid(); arr.fill.fore_color.rgb = OLIVE
    arr.line.fill.background(); arr.shadow.inherit = False

    # APRÈS
    add_round_rect(s, Inches(7.20), Inches(1.85), Inches(5.55), Inches(2.05),
                   WHITE, line_color=OLIVE, line_width=1.5, corner=0.06)
    add_rect(s, Inches(7.20), Inches(1.85), Inches(5.55), Inches(0.40), OLIVE)
    add_text(s, Inches(7.40), Inches(1.85), Inches(5.05), Inches(0.40),
             "PÉRIMÈTRE ÉLARGI — VALIDÉ", size=10, bold=True, color=WHITE,
             anchor=MSO_ANCHOR.MIDDLE)
    add_text(s, Inches(7.40), Inches(2.40), Inches(5.05), Inches(0.50),
             "Pipeline ETL + Plateforme + ML",
             size=17, bold=True, color=NAVY_DARK)
    add_bullets(s, Inches(7.40), Inches(2.95), Inches(5.05), Inches(0.95),
        [
            "Axe 1 — Pipeline ETL multi-sources.",
            "Axe 2 — Plate-forme métier de gouvernance.",
            "Axe 3 — Machine learning multicritère SCAR.",
        ], size=11, bullet_color=OLIVE)

    # 3 axes — triptyque détaillé
    y0 = Inches(4.10)
    h0 = Inches(3.00)
    cw = Inches(4.05)
    gap = Inches(0.10)

    # Axe 1 — Pipeline ETL
    x1 = Inches(0.55)
    add_round_rect(s, x1, y0, cw, h0, NAVY_LIGHT, corner=0.06)
    add_text(s, x1 + Inches(0.30), y0 + Inches(0.20), cw - Inches(0.6), Inches(0.35),
             "AXE 1", size=11, bold=True, color=OLIVE_LIGHT)
    add_text(s, x1 + Inches(0.30), y0 + Inches(0.50), cw - Inches(0.6), Inches(0.55),
             "Pipeline ETL",
             size=16, bold=True, color=WHITE)
    add_bullets(s, x1 + Inches(0.30), y0 + Inches(1.10),
                cw - Inches(0.6), h0 - Inches(1.20),
        [
            "5 sources hétérogènes intégrées.",
            "Imputation hiérarchique 3 niveaux.",
            "Matching flou cédantes (>95 % recall).",
            "Schéma en étoile, traçabilité totale.",
        ], size=10.5, color=WHITE, bullet_color=OLIVE_LIGHT, line_spacing=1.30)

    # Axe 2 — Plate-forme
    x2 = x1 + cw + gap
    add_round_rect(s, x2, y0, cw, h0, NAVY, corner=0.06)
    add_text(s, x2 + Inches(0.30), y0 + Inches(0.20), cw - Inches(0.6), Inches(0.35),
             "AXE 2", size=11, bold=True, color=OLIVE_LIGHT)
    add_text(s, x2 + Inches(0.30), y0 + Inches(0.50), cw - Inches(0.6), Inches(0.55),
             "Plate-forme métier",
             size=16, bold=True, color=WHITE)
    add_bullets(s, x2 + Inches(0.30), y0 + Inches(1.10),
                cw - Inches(0.6), h0 - Inches(1.20),
        [
            "31 pages — dashboard, analyses, pilotage.",
            "Cédantes, courtiers, TTY, FAC, rétro.",
            "Cartographies thématiques + fiche pays.",
            "RBAC + audit d'activité horodaté.",
        ], size=10.5, color=WHITE, bullet_color=OLIVE_LIGHT, line_spacing=1.30)

    # Axe 3 — ML
    x3 = x2 + cw + gap
    add_round_rect(s, x3, y0, cw, h0, OLIVE, corner=0.06)
    add_text(s, x3 + Inches(0.30), y0 + Inches(0.20), cw - Inches(0.6), Inches(0.35),
             "AXE 3", size=11, bold=True, color=NAVY_900)
    add_text(s, x3 + Inches(0.30), y0 + Inches(0.50), cw - Inches(0.6), Inches(0.55),
             "Machine Learning",
             size=16, bold=True, color=WHITE)
    add_bullets(s, x3 + Inches(0.30), y0 + Inches(1.10),
                cw - Inches(0.6), h0 - Inches(1.20),
        [
            "Modèle SCAR — 5 dimensions, score 0-100.",
            "Validation Leave-One-Country-Out (LOCO).",
            "Recommandations Attractif / Neutre / Éviter.",
            "34 pays × 10 ans calibrés.",
        ], size=10.5, color=WHITE, bullet_color=NAVY_DARK, line_spacing=1.30)

    page_footer(s, 13, TOTAL, "03 · Analyse des besoins")


# =============================================================================
#  SECTION 4 — LA SOLUTION
# =============================================================================
def slide_section4_divider():
    section_divider("Section 04",
                    "La solution Atlantic Re",
                    "Architecture trois-tiers et trajectoire de livraison en huit sprints.",
                    14, TOTAL)


def slide_architecture():
    s = prs.slides.add_slide(BLANK)
    page_background(s)
    page_header(s, "04 · Solution Atlantic Re",
                "Architecture trois-tiers — moderne, sécurisée, souveraine")

    # 3 horizontal bands (présentation, métier, données)
    bands = [
        ("PRÉSENTATION", "React 18  ·  TypeScript  ·  Vite  ·  TailwindCSS",
         "31 pages — Axe 1 (Dashboard, Cédantes, Brokers, TTY, FAC, Rétro) + Axe 2 (4 cartographies, fiche pays, comparaison)",
         OLIVE),
        ("MÉTIER",        "FastAPI  ·  Python 3.12  ·  JWT  ·  Pandas",
         "22 routeurs · 25 services · Cœur algorithmique scoring SCAR · Middlewares CORS / RBAC / Rate-limit / Audit",
         NAVY),
        ("DONNÉES",       "Pandas in-memory  ·  MySQL  ·  SQLAlchemy",
         "Excel portefeuille (in-memory) · Base RBAC + journal d'activité · Référentiel pays + 4 tables thématiques externes",
         NAVY_LIGHT),
    ]

    bx = Inches(0.55); bw = Inches(12.20)
    by = Inches(1.95); bh = Inches(1.55); gap = Inches(0.20)
    for i, (head, tech, txt, color) in enumerate(bands):
        y = by + i * (bh + gap)
        add_round_rect(s, bx, y, bw, bh, WHITE,
                       line_color=GRAY_300, line_width=0.5, corner=0.04)
        add_rect(s, bx, y, Inches(0.18), bh, color)
        # left badge
        pill(s, bx + Inches(0.40), y + Inches(0.30), Inches(2.0), Inches(0.40),
             head, color, font_size=11)
        add_text(s, bx + Inches(0.40), y + Inches(0.80), Inches(11.5), Inches(0.35),
                 tech, size=11, bold=True, color=NAVY_DARK)
        add_text(s, bx + Inches(0.40), y + Inches(1.10), Inches(11.5), Inches(0.40),
                 txt, size=10.5, color=GRAY_700)

    # bottom — qualités non fonctionnelles
    qy = Inches(6.95)
    qualities = [
        ("Performance < 2 s", OLIVE),
        ("Sécurité JWT + RBAC", NAVY),
        ("Souveraineté on-prem", NAVY_LIGHT),
        ("Audit horodaté", OLIVE_LIGHT),
        ("Reproductibilité", OLIVE),
    ]
    pw = Inches(2.40); spacing = Inches(0.05)
    total_w = pw * len(qualities) + spacing * (len(qualities) - 1)
    sx = (SW - total_w) / 2
    for i, (txt, c) in enumerate(qualities):
        pill(s, sx + i * (pw + spacing), qy - Inches(0.50),
             pw, Inches(0.40), txt, c, font_size=10)

    page_footer(s, 15, TOTAL, "04 · Solution Atlantic Re")


def slide_roadmap_8_sprints():
    s = prs.slides.add_slide(BLANK)
    page_background(s)
    page_header(s, "04 · Solution Atlantic Re",
                "Trajectoire de livraison — 8 sprints, 5 mois, 4 jalons")

    # Timeline horizontale (S0..S7) — span ~ Inches(12.0)
    tl_x0 = Inches(0.65); tl_y = Inches(2.80); tl_w = Inches(12.05)
    add_rect(s, tl_x0, tl_y + Inches(0.36), tl_w, Inches(0.06), GRAY_300)

    sprints = [
        ("S0", "Immersion\n& cadrage",     "3 sem", NAVY_LIGHT,   None),
        ("S1", "Analyse\n& conception",    "2 sem", NAVY_LIGHT,   "Élargissement"),
        ("S2", "Architecture\n& infra",    "2 sem", NAVY_LIGHT,   None),
        ("S3", "Portail\nDashboard",       "3 sem", OLIVE,         None),
        ("S4", "Portail\nMétier",          "3 sem", OLIVE,         "MVP Axe 1"),
        ("S5", "Données ext.\n& cartos",   "4 sem", OLIVE_LIGHT,   "Axe 2 livré"),
        ("S6", "Modèle\nSCAR",             "2 sem", DANGER,        None),
        ("S7", "Intégration\n& clôture",   "2 sem", NAVY_DARK,     "Soutenance"),
    ]
    n = len(sprints)
    step = tl_w / n
    for i, (lbl, name, dur, c, jalon) in enumerate(sprints):
        cx = tl_x0 + step * i + step/2 - Inches(0.36)
        circle = s.shapes.add_shape(MSO_SHAPE.OVAL, cx, tl_y, Inches(0.72), Inches(0.72))
        circle.fill.solid(); circle.fill.fore_color.rgb = c
        circle.line.color.rgb = WHITE; circle.line.width = Pt(2)
        circle.shadow.inherit = False
        add_text(s, cx, tl_y, Inches(0.72), Inches(0.72), lbl,
                 size=13, bold=True, color=WHITE,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        # name
        add_text(s, cx - Inches(0.30), tl_y + Inches(0.85),
                 Inches(1.30), Inches(0.65),
                 name, size=10, bold=True, color=NAVY_DARK,
                 align=PP_ALIGN.CENTER)
        # durée
        add_text(s, cx - Inches(0.30), tl_y - Inches(0.40),
                 Inches(1.30), Inches(0.30), dur,
                 size=9, italic=True, color=GRAY_500,
                 align=PP_ALIGN.CENTER)
        # jalon
        if jalon:
            mr = s.shapes.add_shape(MSO_SHAPE.DIAMOND,
                                    cx + Inches(0.20), tl_y + Inches(1.55),
                                    Inches(0.30), Inches(0.30))
            mr.fill.solid(); mr.fill.fore_color.rgb = AMBER
            mr.line.color.rgb = NAVY_DARK; mr.line.width = Pt(0.5)
            mr.shadow.inherit = False
            add_text(s, cx - Inches(0.50), tl_y + Inches(1.90),
                     Inches(1.70), Inches(0.30), jalon,
                     size=9, bold=True, color=AMBER, align=PP_ALIGN.CENTER)

    # Phase legend (axis labels)
    add_text(s, Inches(0.65), Inches(2.05), Inches(12), Inches(0.30),
             "FÉVRIER          ·          MARS          ·          AVRIL          ·          MAI          ·          JUIN",
             size=10, bold=True, color=GRAY_500, align=PP_ALIGN.CENTER)

    # bottom — 4 KPIs livraison
    ky = Inches(5.85)
    kpi_card(s, Inches(0.55),  ky, Inches(2.95), Inches(1.20), "31",  "Pages applicatives",  accent=OLIVE)
    kpi_card(s, Inches(3.65),  ky, Inches(2.95), Inches(1.20), "16",  "Besoins fonctionnels", accent=NAVY)
    kpi_card(s, Inches(6.75),  ky, Inches(2.95), Inches(1.20), "34",  "Pays africains couverts", accent=OLIVE_LIGHT)
    kpi_card(s, Inches(9.85),  ky, Inches(2.90), Inches(1.20), "10",  "Exercices historisés", accent=NAVY_LIGHT)

    page_footer(s, 16, TOTAL, "04 · Solution Atlantic Re")


# =============================================================================
#  SECTION 5 — AXE 1 PIPELINE ETL
# =============================================================================
def slide_section5_divider():
    section_divider("Section 05",
                    "Axe 1 — Pipeline ETL\nmulti-sources",
                    "Extraction, transformation, chargement — fondations data du projet.",
                    17, TOTAL)


def slide_etl_architecture():
    s = prs.slides.add_slide(BLANK)
    page_background(s)
    page_header(s, "05 · Axe 1 — Pipeline ETL",
                "Architecture E · T · L — du fichier brut au schéma exploitable")

    # 3 phases en bandes verticales
    phases = [
        ("E — EXTRACT", OLIVE,
         "Collecte multi-sources",
         [
             "5 sources hétérogènes (CSV, XLSX, JSON, API).",
             "Acquisition automatisable Axco / WGI / FMI.",
             "Versionnement par millésime exercice.",
             "Logs de collecte horodatés.",
         ]),
        ("T — TRANSFORM", NAVY,
         "Normalisation & enrichissement",
         [
             "Canonicalisation des codes pays (ISO-3).",
             "Imputation hiérarchique 3 niveaux.",
             "Calcul d'indicateurs dérivés (densité, S/P).",
             "Annotation fiabilité de chaque cellule.",
         ]),
        ("L — LOAD", NAVY_LIGHT,
         "Chargement vers schéma étoile",
         [
             "Tables de référence : ref_pays.",
             "Tables de faits : ext_marche_vie / non_vie.",
             "ext_macroeconomie · ext_gouvernance.",
             "Indexation pays × année pour requêtes < 2 s.",
         ]),
    ]

    cw = Inches(4.05)
    ch = Inches(3.85)
    gap = Inches(0.10)
    y0 = Inches(1.95)
    for i, (head, color, name, items) in enumerate(phases):
        x = Inches(0.55) + i * (cw + gap)
        add_round_rect(s, x, y0, cw, ch, WHITE,
                       line_color=GRAY_300, line_width=0.5, corner=0.06)
        add_rect(s, x, y0, cw, Inches(0.50), color)
        add_text(s, x + Inches(0.20), y0, cw - Inches(0.40), Inches(0.50),
                 head, size=12, bold=True, color=WHITE, anchor=MSO_ANCHOR.MIDDLE)
        add_text(s, x + Inches(0.20), y0 + Inches(0.65), cw - Inches(0.40), Inches(0.50),
                 name, size=14, bold=True, color=NAVY_DARK)
        add_bullets(s, x + Inches(0.20), y0 + Inches(1.20),
                    cw - Inches(0.40), ch - Inches(1.30),
                    items, size=11, bullet_color=color, line_spacing=1.40)

        # arrow between phases
        if i < len(phases) - 1:
            ax = x + cw + Inches(0.005)
            arr = s.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW,
                                     ax - Inches(0.04), y0 + ch/2 - Inches(0.10),
                                     Inches(0.20), Inches(0.20))
            arr.fill.solid(); arr.fill.fore_color.rgb = NAVY
            arr.line.fill.background(); arr.shadow.inherit = False

    # bottom highlight
    add_round_rect(s, Inches(0.55), Inches(6.05), Inches(12.20), Inches(1.00),
                   NAVY_DARK, corner=0.10)
    add_text(s, Inches(0.85), Inches(6.15), Inches(11.7), Inches(0.30),
             "EXIGENCE OPÉRATIONNELLE", size=10, bold=True, color=OLIVE_LIGHT)
    add_text(s, Inches(0.85), Inches(6.45), Inches(11.7), Inches(0.55),
             "Reproductibilité bout-en-bout : tout exercice peut être reconstitué à partir du fichier source\net du script ETL versionné — exigence des comités de risque et d'audit.",
             size=11.5, color=WHITE, italic=True, line_spacing=1.20)

    page_footer(s, 18, TOTAL, "05 · Axe 1 — Pipeline ETL")


def slide_etl_sources():
    s = prs.slides.add_slide(BLANK)
    page_background(s)
    page_header(s, "05 · Axe 1 — Pipeline ETL",
                "Cinq sources de référence — un panel africain inédit")

    rows = [
        ("Axco Navigator", "Marchés Vie / Non-Vie",
         "Primes, pénétration, densité, croissance, S/P par marché",
         "Référence sectorielle"),
        ("WGI Banque Mondiale", "Gouvernance",
         "6 indicateurs de qualité institutionnelle (1996-2024)",
         "Standard académique"),
        ("FMI WEO", "Macroéconomie",
         "PIB, croissance, inflation, PIB/hab PPP, balance courante",
         "Source officielle"),
        ("FANAF", "Marché francophone",
         "Statistiques zone CIMA — primes, sinistres, ratios",
         "Réseau institutionnel"),
        ("Atlas Magazine / Sigma", "Concurrence",
         "Top réassureurs, parts de marché, classements continentaux",
         "Veille sectorielle"),
    ]

    nx = Inches(0.55)
    ny = Inches(1.95)
    cwd = [Inches(2.85), Inches(2.40), Inches(4.55), Inches(2.40)]
    rh = Inches(0.55)

    headers = ["Source", "Périmètre", "Variables clés", "Justification"]
    cx = nx
    for i, w in enumerate(cwd):
        add_rect(s, cx, ny, w, rh, NAVY)
        add_text(s, cx + Inches(0.15), ny, w - Inches(0.30), rh,
                 headers[i], size=11, bold=True, color=WHITE,
                 anchor=MSO_ANCHOR.MIDDLE)
        cx += w
    for r, row in enumerate(rows, start=1):
        cx = nx
        bg = WHITE if r % 2 == 1 else SURFACE_2
        for i, w in enumerate(cwd):
            add_rect(s, cx, ny + rh * r, w, rh, bg)
            color = NAVY_DARK
            bold = (i == 0)
            if i == 1:
                color = OLIVE
                bold = True
            elif i == 3:
                color = GRAY_700
            add_text(s, cx + Inches(0.15), ny + rh * r, w - Inches(0.30), rh,
                     row[i], size=10.5, bold=bold, color=color,
                     anchor=MSO_ANCHOR.MIDDLE)
            cx += w

    # KPIs en bas
    ky = Inches(5.55)
    kpi_card(s, Inches(0.55), ky, Inches(2.95), Inches(1.20),
             "34", "Pays africains intégrés", accent=OLIVE)
    kpi_card(s, Inches(3.65), ky, Inches(2.95), Inches(1.20),
             "10 ans", "Profondeur historique (2015-2024)", accent=NAVY)
    kpi_card(s, Inches(6.75), ky, Inches(2.95), Inches(1.20),
             "> 95 %", "Recall matching cédantes", accent=OLIVE_LIGHT)
    kpi_card(s, Inches(9.85), ky, Inches(2.90), Inches(1.20),
             "5", "Sources fédérées", accent=NAVY_LIGHT)

    # bandeau
    add_round_rect(s, Inches(0.55), Inches(6.95), Inches(12.20), Inches(0.18),
                   OLIVE, corner=0.30)

    page_footer(s, 19, TOTAL, "05 · Axe 1 — Pipeline ETL")


def slide_etl_imputation():
    s = prs.slides.add_slide(BLANK)
    page_background(s)
    page_header(s, "05 · Axe 1 — Pipeline ETL",
                "Imputation hiérarchique — fiabiliser sans inventer")

    # KPI haut : volumétrie
    kx = Inches(0.55); ky = Inches(1.95); kw = Inches(2.95)
    cards = [
        ("12 / 34",  "Pays à données complètes",      OLIVE),
        ("22 / 34",  "Pays nécessitant imputation",   WARNING),
        ("3",        "Niveaux de fallback",           NAVY),
        ("100 %",    "Cellules annotées fiabilité",   OLIVE_LIGHT),
    ]
    for i, (v, l, c) in enumerate(cards):
        x = kx + i * (kw + Inches(0.10))
        add_round_rect(s, x, ky, kw, Inches(1.10), c, corner=0.08)
        add_text(s, x, ky + Inches(0.08), kw, Inches(0.55),
                 v, size=22, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        add_text(s, x, ky + Inches(0.65), kw, Inches(0.40),
                 l, size=10.5, color=WHITE, align=PP_ALIGN.CENTER)

    # 3 niveaux d'imputation
    levels = [
        ("NIVEAU 1", "Interpolation temporelle",
         "Linéaire ou spline cubique sur la chronique du même pays.",
         "Fiabilité élevée — comble les trous courts.", OLIVE),
        ("NIVEAU 2", "Inférence régionale",
         "Médiane par bloc régional (CIMA · SADC · CEDEAO).",
         "Fiabilité moyenne — capture l'effet zone monétaire.", AMBER),
        ("NIVEAU 3", "Apprentissage supervisé",
         "Régression à partir de variables corrélées (PIB, démographie).",
         "Fiabilité contrôlée — annotée pour transparence.", NAVY_LIGHT),
    ]
    by = Inches(3.40); bw = Inches(3.95); bh = Inches(2.50)
    for i, (head, name, desc, qual, color) in enumerate(levels):
        x = Inches(0.55) + i * (bw + Inches(0.15))
        add_round_rect(s, x, by, bw, bh, WHITE,
                       line_color=GRAY_300, line_width=0.5, corner=0.06)
        add_rect(s, x, by, bw, Inches(0.40), color)
        add_text(s, x, by, bw, Inches(0.40),
                 head, size=11, bold=True, color=WHITE,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        add_text(s, x + Inches(0.20), by + Inches(0.55), bw - Inches(0.40),
                 Inches(0.45),
                 name, size=14, bold=True, color=NAVY_DARK)
        add_text(s, x + Inches(0.20), by + Inches(1.10), bw - Inches(0.40),
                 Inches(0.85),
                 desc, size=11, color=GRAY_700, line_spacing=1.30)
        add_round_rect(s, x + Inches(0.20), by + bh - Inches(0.65),
                       bw - Inches(0.40), Inches(0.50),
                       SURFACE_2, corner=0.20)
        add_text(s, x + Inches(0.20), by + bh - Inches(0.65),
                 bw - Inches(0.40), Inches(0.50),
                 qual, size=9.5, italic=True, color=NAVY_DARK,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

    page_footer(s, 20, TOTAL, "05 · Axe 1 — Pipeline ETL")


# =============================================================================
#  SECTION 6 — AXE 2 PLATE-FORME MÉTIER
# =============================================================================
def slide_section6_plateforme_divider():
    section_divider("Section 06",
                    "Axe 2 — Plate-forme\nmétier de gouvernance",
                    "Restitution, analyses, pilotage et gouvernance — résultats par module.",
                    21, TOTAL)


def _slide_axe_module(idx_page, eyebrow, title, kpis, bullets,
                      shot_label, section_lbl):
    """Layout générique : KPI cards en haut + bullets + zone capture droite."""
    s = prs.slides.add_slide(BLANK)
    page_background(s)
    page_header(s, eyebrow, title)

    # KPI horizontal strip
    if kpis:
        kx = Inches(0.55); ky = Inches(1.95)
        kw = (Inches(5.5) - Inches(0.20) * (len(kpis) - 1)) / len(kpis) if len(kpis) > 1 else Inches(5.5)
        # use 3 stacked compact KPIs if 3
        for i, (val, lbl) in enumerate(kpis):
            x = kx + i * (kw + Inches(0.20))
            add_round_rect(s, x, ky, kw, Inches(1.10),
                           NAVY, corner=0.08)
            add_text(s, x, ky + Inches(0.08), kw, Inches(0.55),
                     val, size=22, bold=True, color=OLIVE_LIGHT,
                     align=PP_ALIGN.CENTER)
            add_text(s, x, ky + Inches(0.65), kw, Inches(0.40),
                     lbl, size=10, color=WHITE,
                     align=PP_ALIGN.CENTER)

    # Bullets fonctionnalités
    by = Inches(3.30) if kpis else Inches(2.05)
    add_text(s, Inches(0.55), by, Inches(5.7), Inches(0.40),
             "Fonctionnalités livrées", size=13, bold=True, color=NAVY_DARK)
    add_bullets(s, Inches(0.55), by + Inches(0.45),
                Inches(5.7), Inches(3.0),
                bullets, size=12, line_spacing=1.45)

    # Capture zone droite
    add_screenshot_zone(s, Inches(6.50), Inches(1.95),
                        Inches(6.30), Inches(5.10),
                        label=shot_label)

    page_footer(s, idx_page, TOTAL, section_lbl)


def slide_axe1_dashboard():
    _slide_axe_module(
        22, "06 · Axe 2 — Restitution",
        "Tableau de bord consolidé du portefeuille",
        kpis=[("4", "KPIs d'en-tête"), ("6", "Vues croisées"), ("Temps réel", "Mise à jour")],
        bullets=[
            "Primes écrites, résultat technique, ULR pondéré, nombre de contrats.",
            "Évolution temporelle pluriannuelle + alertes ULR critiques.",
            "Répartitions par branche, par origine FAC vs Traités.",
            "Cartographie mondiale interactive avec tooltips.",
        ],
        shot_label="Capture — Dashboard principal (route /dashboard)",
        section_lbl="06 · Axe 2 — Plate-forme métier")


def slide_axe1_filtrage():
    _slide_axe_module(
        23, "06 · Axe 2 — Restitution",
        "Filtrage multicritère — exploration fine du portefeuille",
        kpis=[("15+", "Critères croisables"), ("3", "Niveaux de filtres"), ("< 2 s", "Restitution")],
        bullets=[
            "Filtres identitaires : exercice, cédante, courtier, statut.",
            "Filtres d'analyse : branche, pays, type de contrat (FAC / TTY / TTE).",
            "Filtres financiers : seuils de primes, ULR, parts souscrites.",
            "État sérialisable — partage de vue d'analyse entre utilisateurs.",
        ],
        shot_label="Capture — Panneau de filtrage global déployé",
        section_lbl="06 · Axe 2 — Plate-forme métier")


def slide_axe1_cedantes_courtiers():
    _slide_axe_module(
        24, "06 · Axe 2 — Analyses",
        "Analyses cédantes & courtiers — réconciliation par matching flou",
        kpis=[("> 95 %", "Recall matching"), ("RapidFuzz", "Levenshtein"),
              ("Jellyfish", "Phonétique")],
        bullets=[
            "Profil cédante : historique, ULR, indice de diversification, répartition par branche.",
            "Matching flou des dénominations — RapidFuzz (Levenshtein) + Jellyfish (phonétique).",
            "Réconciliation automatique de jusqu'à 6 variantes par cédante.",
            "Profil courtier : volume intermédié, branches, cédantes apportées.",
            "Forage agrégé → contrats élémentaires en un clic.",
        ],
        shot_label="Capture — Fiche cédante (route /analyse-cedante/:cedante)",
        section_lbl="06 · Axe 2 — Plate-forme métier")


def slide_axe1_comparaison_tty():
    _slide_axe_module(
        25, "06 · Axe 2 — Pilotage",
        "Comparaison directe & cibles TTY",
        kpis=[("2-4", "Entités comparables"), ("5", "Axes radar"), ("3", "Filtres rapides")],
        bullets=[
            "Comparaison côte-à-côte : Pays × Branche, Cédante vs Cédante.",
            "Radar normalisé sur 5 axes (Primes, Résultat, ULR, Diversité, SI).",
            "Cibles TTY : repérage des cédantes à transformer en traité additionnel.",
            "Boutons Stable / Baisse / Hausse — lecture instantanée du potentiel.",
        ],
        shot_label="Capture — Module Comparaison ou Cibles TTY",
        section_lbl="06 · Axe 2 — Plate-forme métier")


def slide_axe1_fac_retro():
    _slide_axe_module(
        26, "06 · Axe 2 — Pilotage",
        "Saturation FAC & Rétrocession",
        kpis=[("Auto", "Détection abus FAC"), ("EPI / PMD", "Calculs métier"),
              ("Multi", "Partenaires")],
        bullets=[
            "Algorithme OR — détection automatique des cédantes en sur-utilisation FAC.",
            "Surbrillance rouge des branches au-delà des seuils métier.",
            "Module rétrocession : EPI, PMD, Taux de placement.",
            "Vue dédiée Sécurités — analyse des partenaires porteurs.",
        ],
        shot_label="Capture — FAC Saturation ou Rétrocession Affaires Traitées",
        section_lbl="06 · Axe 2 — Plate-forme métier")


def slide_axe1_admin():
    _slide_axe_module(
        27, "06 · Axe 2 — Gouvernance",
        "Administration, sécurité et auditabilité",
        kpis=[("3", "Profils RBAC"), ("100 %", "Actions tracées"), ("JWT", "Authentification")],
        bullets=[
            "RBAC — Administrateur · Souscripteur · Lecteur.",
            "CRUD complet utilisateurs + réinitialisation mot de passe.",
            "Journal d'activité horodaté et immuable.",
            "Configuration des sources Excel persistées sans reload backend.",
        ],
        shot_label="Capture — Panneau Administration (route /admin)",
        section_lbl="06 · Axe 2 — Plate-forme métier")


# =============================================================================
#  SECTION 7 — AXE 3 MACHINE LEARNING
# =============================================================================
def slide_section7_ml_divider():
    section_divider("Section 07",
                    "Axe 3 — Machine Learning\nmulticritère SCAR",
                    "Cartographies thématiques, fiches pays, scoring SCAR et recommandations.",
                    28, TOTAL)


def slide_axe2_panel():
    s = prs.slides.add_slide(BLANK)
    page_background(s)
    page_header(s, "07 · Axe 3 — Données du modèle",
                "Le panel Africa 2030 — fondations du scoring ML")

    # 4 KPI strip
    kx = Inches(0.55); ky = Inches(1.95); kw = Inches(2.95)
    cards = [
        ("34",   "Pays africains",       OLIVE),
        ("10",   "Exercices (2015-24)",  NAVY),
        ("4",    "Tables thématiques",   OLIVE_LIGHT),
        ("3",    "Niveaux d'imputation", NAVY_LIGHT),
    ]
    for i, (v, l, c) in enumerate(cards):
        x = kx + i * (kw + Inches(0.10))
        add_round_rect(s, x, ky, kw, Inches(1.10), c, corner=0.08)
        add_text(s, x, ky + Inches(0.08), kw, Inches(0.55),
                 v, size=24, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        add_text(s, x, ky + Inches(0.62), kw, Inches(0.40),
                 l, size=11, color=WHITE, align=PP_ALIGN.CENTER)

    # 4 thématiques cards
    by = Inches(3.40); bw = Inches(2.95); bh = Inches(2.30)
    themes = [
        ("MARCHÉ NON-VIE", "Primes, pénétration,\ndensité, croissance, S/P", "Axco Navigator"),
        ("MARCHÉ VIE",     "Primes, pénétration,\ndensité, croissance",      "Axco Navigator"),
        ("MACROÉCONOMIE",  "PIB, croissance, inflation,\nPIB/hab PPP, courant", "FMI WEO"),
        ("GOUVERNANCE",    "Stabilité, qualité régul.,\nKaopen, IDE",        "BM WGI · Chinn-Ito"),
    ]
    for i, (h, txt, src) in enumerate(themes):
        x = kx + i * (bw + Inches(0.10))
        add_round_rect(s, x, by, bw, bh, WHITE,
                       line_color=GRAY_300, line_width=0.5, corner=0.06)
        add_rect(s, x, by, bw, Inches(0.40), NAVY)
        add_text(s, x, by, bw, Inches(0.40),
                 h, size=10, bold=True, color=OLIVE_LIGHT,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        add_text(s, x + Inches(0.20), by + Inches(0.55), bw - Inches(0.40),
                 Inches(1.40),
                 txt, size=12, color=NAVY_DARK)
        add_round_rect(s, x + Inches(0.20), by + bh - Inches(0.50),
                       bw - Inches(0.40), Inches(0.35), SURFACE_2, corner=0.30)
        add_text(s, x + Inches(0.20), by + bh - Inches(0.50),
                 bw - Inches(0.40), Inches(0.35),
                 src, size=10, italic=True, color=GRAY_700,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

    # Bottom highlight
    add_round_rect(s, Inches(0.55), Inches(6.10), Inches(12.20), Inches(0.95),
                   NAVY_DARK, corner=0.10)
    add_text(s, Inches(0.85), Inches(6.20), Inches(11.7), Inches(0.30),
             "DIFFÉRENCIATEUR", size=10, bold=True, color=OLIVE_LIGHT)
    add_text(s, Inches(0.85), Inches(6.50), Inches(11.7), Inches(0.50),
             "Chaque valeur imputée est annotée par sa méthode d'imputation et son niveau de fiabilité —\ntraçabilité totale exigée par les comités de risque.",
             size=11.5, color=WHITE, italic=True, line_spacing=1.20)

    page_footer(s, 29, TOTAL, "07 · Axe 3 — Machine Learning")


def slide_axe2_carto():
    _slide_axe_module(
        30, "07 · Axe 3 — Restitution",
        "Cartographies thématiques interactives — 4 dimensions",
        kpis=[("4", "Modules cartos"), ("Top 10", "Classement"), ("Choro.", "Carte")],
        bullets=[
            "Carte choroplèthe Afrique avec sélecteur indicateur × année.",
            "Indicateurs synthétiques + Top 10 dynamique.",
            "Vue croisée configurable (deux indicateurs simultanés).",
            "Séries temporelles, matrice pays × années, distributions, classement triable.",
        ],
        shot_label="Capture — Cartographie Non-Vie ou Macroéconomie",
        section_lbl="07 · Axe 3 — Machine Learning")


def slide_axe2_fiche_pays():
    _slide_axe_module(
        31, "07 · Axe 3 — Restitution",
        "Fiche pays multidimensionnelle & comparaison",
        kpis=[("4", "Dimensions"), ("Radar", "Profil"), ("2", "Pays comparables")],
        bullets=[
            "Profil consolidé : marché Vie + Non-Vie + macro + gouvernance.",
            "Graphique radar et positionnement continental relatif.",
            "Comparaison côte-à-côte de deux pays sur l'ensemble des dimensions.",
            "Commentaires analytiques générés dynamiquement à partir des données.",
        ],
        shot_label="Capture — Fiche pays ou Comparaison transnationale",
        section_lbl="07 · Axe 3 — Machine Learning")


def slide_axe2_scar():
    s = prs.slides.add_slide(BLANK)
    page_background(s)
    page_header(s, "07 · Axe 3 — Scoring SCAR",
                "Le modèle SCAR — synthèse multicritère pour la décision")

    # Texte gauche
    add_text(s, Inches(0.55), Inches(1.95), Inches(6.0), Inches(0.40),
             "Cinq dimensions agrégées en un score 0-100",
             size=14, bold=True, color=NAVY_DARK)

    dims = [
        ("Économique",       "PIB, croissance, inflation, PIB/hab PPP, intégration", OLIVE),
        ("Marché Assurance", "Primes, pénétration, densité, S/P",                    OLIVE_LIGHT),
        ("Réglementaire",    "WGI, stabilité politique, Kaopen, IDE",                NAVY),
        ("Concurrentielle",  "Nb réassureurs, part de marché Top 3",                 NAVY_LIGHT),
        ("Interne",          "Primes nettes, ULR, résultat, sinistres",              DANGER),
    ]
    y0 = Inches(2.45)
    for i, (n, d, c) in enumerate(dims):
        y = y0 + i * Inches(0.62)
        add_round_rect(s, Inches(0.55), y, Inches(6.20), Inches(0.55),
                       WHITE, line_color=GRAY_300, line_width=0.5, corner=0.10)
        add_rect(s, Inches(0.55), y, Inches(0.10), Inches(0.55), c)
        add_text(s, Inches(0.85), y, Inches(2.0), Inches(0.55),
                 n, size=12, bold=True, color=NAVY_DARK,
                 anchor=MSO_ANCHOR.MIDDLE)
        add_text(s, Inches(2.85), y, Inches(3.85), Inches(0.55),
                 d, size=10.5, color=GRAY_700,
                 anchor=MSO_ANCHOR.MIDDLE)

    # Droite : badges sortie + état
    add_text(s, Inches(7.00), Inches(1.95), Inches(5.8), Inches(0.40),
             "Sorties opérationnelles", size=14, bold=True, color=NAVY_DARK)

    badges = [
        ("ATTRACTIF", OLIVE,    "Marché à prioriser — souscription active"),
        ("NEUTRE",    AMBER,    "Maintenir — surveiller les inflexions"),
        ("À ÉVITER",  DANGER,   "Risque > rentabilité — désengagement"),
    ]
    for i, (lbl, c, desc) in enumerate(badges):
        y = Inches(2.45) + i * Inches(0.95)
        add_round_rect(s, Inches(7.00), y, Inches(5.75), Inches(0.80),
                       WHITE, line_color=GRAY_300, line_width=0.5, corner=0.06)
        pill(s, Inches(7.20), y + Inches(0.20), Inches(1.55), Inches(0.40),
             lbl, c, font_size=11)
        add_text(s, Inches(8.95), y + Inches(0.10), Inches(3.7), Inches(0.65),
                 desc, size=11, color=GRAY_700,
                 anchor=MSO_ANCHOR.MIDDLE)

    # Status pill
    add_round_rect(s, Inches(7.00), Inches(5.45), Inches(5.75), Inches(1.55),
                   NAVY, corner=0.06)
    add_text(s, Inches(7.20), Inches(5.55), Inches(5.4), Inches(0.30),
             "ÉTAT D'AVANCEMENT", size=10, bold=True, color=OLIVE_LIGHT)
    add_text(s, Inches(7.20), Inches(5.85), Inches(5.4), Inches(0.45),
             "Composante en cours de finalisation",
             size=14, bold=True, color=WHITE)
    add_text(s, Inches(7.20), Inches(6.32), Inches(5.4), Inches(0.55),
             "Calibration et validation Leave-One-Country-Out\nlivrées au Sprint 6 (juin 2026).",
             size=10.5, color=GRAY_300, italic=True, line_spacing=1.20)

    page_footer(s, 32, TOTAL, "07 · Axe 3 — Machine Learning")


def slide_ml_pipeline_recommendations():
    s = prs.slides.add_slide(BLANK)
    page_background(s)
    page_header(s, "07 · Axe 3 — Pipeline ML",
                "Entraînement, validation LOCO et recommandations actionnables")

    # Pipeline ML horizontal
    add_text(s, Inches(0.55), Inches(1.85), Inches(12.2), Inches(0.40),
             "Pipeline d'apprentissage et de scoring",
             size=14, bold=True, color=NAVY_DARK)

    steps = [
        ("01", "Préparation",
         "Standardisation\nFeature engineering",  OLIVE),
        ("02", "Entraînement",
         "scikit-learn · pondération\ndes 5 dimensions SCAR", NAVY),
        ("03", "Validation",
         "Leave-One-Country-Out\n(34 itérations)",   NAVY_LIGHT),
        ("04", "Scoring",
         "Score 0-100\nattribué par marché",        OLIVE_LIGHT),
        ("05", "Recommandation",
         "Attractif / Neutre /\nÀ éviter",          DANGER),
    ]
    n = len(steps)
    avail = Inches(12.20)
    sw = avail / n
    yT = Inches(2.45)
    for i, (num, name, desc, c) in enumerate(steps):
        x = Inches(0.55) + sw * i
        # circle
        cx = x + sw/2 - Inches(0.34)
        circle = s.shapes.add_shape(MSO_SHAPE.OVAL, cx, yT, Inches(0.68), Inches(0.68))
        circle.fill.solid(); circle.fill.fore_color.rgb = c
        circle.line.color.rgb = WHITE; circle.line.width = Pt(2)
        circle.shadow.inherit = False
        add_text(s, cx, yT, Inches(0.68), Inches(0.68), num,
                 size=14, bold=True, color=WHITE,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        add_text(s, x, yT + Inches(0.78), sw, Inches(0.30),
                 name, size=11, bold=True, color=NAVY_DARK,
                 align=PP_ALIGN.CENTER)
        add_text(s, x, yT + Inches(1.10), sw, Inches(0.85),
                 desc, size=10, color=GRAY_700,
                 align=PP_ALIGN.CENTER, line_spacing=1.30)
        if i < n-1:
            add_rect(s, cx + Inches(0.68), yT + Inches(0.31),
                     sw - Inches(0.68), Inches(0.06), GRAY_300)

    # Trois recommandations finales
    rby = Inches(4.55)
    cards = [
        ("ATTRACTIF", OLIVE,
         "Score ≥ 70",
         "Marché à prioriser — souscription active, allocation de capacité renforcée."),
        ("NEUTRE", AMBER,
         "Score 40-69",
         "Maintenir position — surveiller les inflexions macro et réglementaires."),
        ("À ÉVITER", DANGER,
         "Score < 40",
         "Risque > rentabilité — désengagement progressif ou capacité plafonnée."),
    ]
    cw = Inches(4.05); ch = Inches(2.00); gap = Inches(0.10)
    for i, (lbl, color, score, desc) in enumerate(cards):
        x = Inches(0.55) + i * (cw + gap)
        add_round_rect(s, x, rby, cw, ch, WHITE,
                       line_color=GRAY_300, line_width=0.5, corner=0.06)
        add_rect(s, x, rby, cw, Inches(0.45), color)
        add_text(s, x, rby, cw, Inches(0.45),
                 lbl, size=12, bold=True, color=WHITE,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        add_text(s, x + Inches(0.20), rby + Inches(0.55), cw - Inches(0.40),
                 Inches(0.45),
                 score, size=14, bold=True, color=NAVY_DARK,
                 align=PP_ALIGN.CENTER)
        add_text(s, x + Inches(0.20), rby + Inches(1.00), cw - Inches(0.40),
                 Inches(0.95),
                 desc, size=10.5, color=GRAY_700,
                 align=PP_ALIGN.CENTER, line_spacing=1.30)

    # Bottom — message clé pour la DG
    add_round_rect(s, Inches(0.55), Inches(6.70), Inches(12.20), Inches(0.40),
                   NAVY_DARK, corner=0.20)
    add_text(s, Inches(0.85), Inches(6.70), Inches(11.7), Inches(0.40),
             "Validation LOCO = robustesse hors-échantillon — chaque pays scoré sans avoir vu son historique d'entraînement.",
             size=10.5, italic=True, color=OLIVE_LIGHT,
             anchor=MSO_ANCHOR.MIDDLE)

    page_footer(s, 33, TOTAL, "07 · Axe 3 — Machine Learning")


# =============================================================================
#  SECTION 8 — CONVERGENCE & BILAN
# =============================================================================
def slide_section8_divider():
    section_divider("Section 08",
                    "Convergence & bilan",
                    "Croisement interne / externe, indicateurs de réalisation et perspectives.",
                    34, TOTAL)


def slide_convergence():
    s = prs.slides.add_slide(BLANK)
    page_background(s)
    page_header(s, "08 · Convergence & bilan",
                "Vue de synergie — identifier les marchés sous-exploités")

    # gauche : explication
    add_text(s, Inches(0.55), Inches(1.95), Inches(5.7), Inches(0.40),
             "Croiser exposition propre × attractivité de marché",
             size=14, bold=True, color=NAVY_DARK)
    add_bullets(s, Inches(0.55), Inches(2.45), Inches(5.7), Inches(2.5),
        [
            "Mise en regard : exposition propre (Axe 2) × attractivité de marché (Axe 3).",
            "Calcul d'un taux de pénétration relatif par marché.",
            "Marchés attractifs mais sous-exploités → cibles de développement.",
            "Marchés sur-exposés → arbitrages d'allocation à instruire.",
        ], size=12)

    # 3 quadrant labels
    qy = Inches(5.10)
    quads = [
        ("Sous-exploités\nattractifs",   OLIVE,        "Cible prioritaire"),
        ("Bien positionnés",             NAVY_LIGHT,   "Maintenir"),
        ("Sur-exposés\nrisqués",         DANGER,       "Arbitrage"),
    ]
    qw = Inches(1.85); gap = Inches(0.05)
    sx = Inches(0.55)
    for i, (n, c, tag) in enumerate(quads):
        x = sx + i * (qw + gap)
        add_round_rect(s, x, qy, qw, Inches(1.55),
                       WHITE, line_color=GRAY_300, line_width=0.5, corner=0.08)
        add_rect(s, x, qy, qw, Inches(0.40), c)
        add_text(s, x, qy + Inches(0.45), qw, Inches(0.65),
                 n, size=11, bold=True, color=NAVY_DARK,
                 align=PP_ALIGN.CENTER)
        add_text(s, x, qy + Inches(1.10), qw, Inches(0.40),
                 tag, size=10, italic=True, color=GRAY_500,
                 align=PP_ALIGN.CENTER)

    # capture droite
    add_screenshot_zone(s, Inches(6.50), Inches(1.95),
                        Inches(6.30), Inches(5.10),
                        label="Capture — Page Synergie (route /analyse-synergie)")

    page_footer(s, 35, TOTAL, "08 · Convergence & bilan")


# =============================================================================
#  BILAN ET CLÔTURE (toujours dans la section 08)
# =============================================================================
def slide_bilan_resultats():
    s = prs.slides.add_slide(BLANK)
    page_background(s)
    page_header(s, "08 · Convergence & bilan",
                "Indicateurs de réalisation — ce qui a été livré")

    # Big KPI grid — récap des 3 axes
    kpis = [
        ("5",      "Sources ETL fédérées",        OLIVE),
        ("> 95 %", "Recall matching cédantes",    NAVY),
        ("31",     "Pages applicatives livrées",  NAVY_LIGHT),
        ("16",     "Besoins fonctionnels",        OLIVE_LIGHT),
        ("34",     "Pays scorés par SCAR",        OLIVE),
        ("5",      "Dimensions du modèle ML",     NAVY),
        ("LOCO",   "Validation hors-échantillon", NAVY_LIGHT),
        ("100 %",  "Actions tracées + RBAC",      OLIVE_LIGHT),
    ]
    cols = 4
    cw = Inches(2.95); ch = Inches(1.45)
    gx = Inches(0.10); gy = Inches(0.20)
    sx = Inches(0.55); sy = Inches(1.95)
    for i, (v, l, c) in enumerate(kpis):
        col = i % cols; row = i // cols
        x = sx + col * (cw + gx)
        y = sy + row * (ch + gy)
        kpi_card(s, x, y, cw, ch, v, l, accent=c)

    # bottom banner with key business outcomes
    add_round_rect(s, Inches(0.55), Inches(5.40), Inches(12.20), Inches(1.65),
                   NAVY_DARK, corner=0.06)
    add_text(s, Inches(0.85), Inches(5.50), Inches(11.7), Inches(0.35),
             "POUR LA DIRECTION GÉNÉRALE", size=11, bold=True, color=OLIVE_LIGHT)
    add_text(s, Inches(0.85), Inches(5.85), Inches(11.7), Inches(0.50),
             "Trois axes — un pipeline ETL, une plate-forme métier, un modèle ML — au service de la décision.",
             size=14, bold=True, color=WHITE)
    add_bullets(s, Inches(0.85), Inches(6.35), Inches(11.7), Inches(0.75),
        [
            "Axe 1 ETL — fondations data fiables, traçables, reproductibles.",
            "Axe 2 Plate-forme — fiabilisation cédantes + temps d'analyse réduit (jours → minutes).",
            "Axe 3 ML — cadre quantitatif SCAR pour objectiver les arbitrages Reach2030.",
        ], size=11, color=WHITE, bullet_color=OLIVE_LIGHT, line_spacing=1.20)

    page_footer(s, 36, TOTAL, "08 · Convergence & bilan")


# =============================================================================
#  CONSTRUCTION
# =============================================================================
slide_cover()                          #  1
slide_agenda()                         #  2
slide_section1_divider()               #  3
slide_reach2030()                      #  4
slide_pole_bd()                        #  5
slide_section2_divider()               #  6
slide_marche_africain()                #  7
slide_paysage_concurrence()            #  8
slide_benchmark_outils()               #  9
slide_section3_divider()               # 10
slide_demarche_3_cycles()              # 11
slide_constats_terrain()               # 12
slide_elargissement_perimetre()        # 13
slide_section4_divider()               # 14
slide_architecture()                   # 15
slide_roadmap_8_sprints()              # 16
slide_section5_divider()               # 17  Axe 1 — Pipeline ETL
slide_etl_architecture()               # 18
slide_etl_sources()                    # 19
slide_etl_imputation()                 # 20
slide_section6_plateforme_divider()    # 21  Axe 2 — Plate-forme métier
slide_axe1_dashboard()                 # 22
slide_axe1_filtrage()                  # 23
slide_axe1_cedantes_courtiers()        # 24
slide_axe1_comparaison_tty()           # 25
slide_axe1_fac_retro()                 # 26
slide_axe1_admin()                     # 27
slide_section7_ml_divider()            # 28  Axe 3 — Machine Learning
slide_axe2_panel()                     # 29
slide_axe2_carto()                     # 30
slide_axe2_fiche_pays()                # 31
slide_axe2_scar()                      # 32
slide_ml_pipeline_recommendations()    # 33
slide_section8_divider()               # 34  Convergence & bilan
slide_convergence()                    # 35
slide_bilan_resultats()                # 36


out = "Presentation_Atlantic_Re_DG.pptx"
prs.save(out)
print(f"OK - {len(prs.slides)} slides -> {out}")
